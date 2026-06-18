from __future__ import annotations

from pathlib import Path

from finalreview.models import DocumentChunk
from finalreview.parsers.base import BaseParser
from finalreview.utils.text_utils import normalize_text


class PPTXParser(BaseParser):
    file_types = {".pptx"}

    def parse(self, file_path: Path) -> list[DocumentChunk]:
        try:
            from pptx import Presentation
        except Exception as exc:
            return [
                self.make_chunk(
                    file_path,
                    f"PPTX 解析依赖 python-pptx 不可用：{exc}",
                    0,
                    metadata={"parse_error": "missing_python_pptx"},
                )
            ]
        chunks: list[DocumentChunk] = []
        try:
            prs = Presentation(str(file_path))
            for idx, slide in enumerate(prs.slides, start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        texts.append(shape.text)
                notes = self._extract_notes(slide)
                if notes:
                    texts.append(f"备注：{notes}")
                text = normalize_text("\n".join(texts)) or "该幻灯片未提取到文本。"
                chunks.append(self.make_chunk(file_path, text, idx, slide_number=idx))
        except Exception as exc:
            chunks.append(
                self.make_chunk(
                    file_path,
                    f"PPTX 解析失败：{exc}",
                    0,
                    metadata={"parse_error": str(exc)},
                )
            )
        return chunks

    def _extract_notes(self, slide) -> str:
        try:
            notes_slide = slide.notes_slide
            return "\n".join(shape.text for shape in notes_slide.shapes if hasattr(shape, "text"))
        except Exception:
            return ""
